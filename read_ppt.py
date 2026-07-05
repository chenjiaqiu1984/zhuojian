from pptx import Presentation
prs = Presentation(r'K:\Code_folder\ohcard\卓见心理·艺术疗愈.pptx')
for i, slide in enumerate(prs.slides):
    print(f'=={i+1}==')
    for shape in slide.shapes:
        if shape.has_text_frame:
            print(shape.text_frame.text)
